import os
import tkinter as tk
from tkinter import filedialog, ttk, messagebox
import threading
import random
from pptx import Presentation
import PyPDF2
import fitz  # PyMuPDF for better PDF extraction
import genanki
import anthropic
import json
import re
import time

class ClaudeEnhancedConverter:
    def __init__(self, api_key):
        # Initialize Claude API client
        self.client = anthropic.Anthropic(api_key=api_key)
        
        # Create a unique model ID for Anki
        self.model_id = random.randrange(1 << 30, 1 << 31)
        # Define the card template
        self.model = genanki.Model(
            self.model_id,
            'Claude-Enhanced Presentation Card',
            fields=[
                {'name': 'Question'},
                {'name': 'Answer'},
                {'name': 'Slide'},
                {'name': 'Context'},
            ],
            templates=[
                {
                    'name': 'Card',
                    'qfmt': '{{Question}}',
                    'afmt': '{{FrontSide}}<hr id="answer">{{Answer}}<br><br><i>Slide: {{Slide}}</i>',
                },
            ])
    
    def extract_from_pptx(self, pptx_path, progress_callback=None):
        """Extract text from PowerPoint presentation"""
        slides_content = []
        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)
        
        for i, slide in enumerate(prs.slides, 1):
            if progress_callback:
                progress_callback(10 + (i / total_slides * 20), f"Extracting slide {i}/{total_slides}...")
                
            slide_title = ""
            slide_content = ""
            
            # Extract slide title and content
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        if not slide_title and shape.text.strip():
                            slide_title = shape.text.strip()
                        else:
                            slide_content += shape.text.strip() + "\n"
            
            # Store slide info even if title or content is empty
            slides_content.append({
                'title': slide_title,
                'content': slide_content,
                'slide_num': i
            })
            
            # Debug info
            print(f"Extracted PPTX Slide {i}:")
            print(f"  Title: {slide_title}")
            print(f"  Content length: {len(slide_content)}")
        
        return slides_content
    
    def extract_from_pdf(self, pdf_path, progress_callback=None):
        """Extract text from PDF presentation using PyMuPDF for better extraction"""
        slides_content = []
        
        # Try PyMuPDF first (better extraction)
        try:
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            
            for i, page in enumerate(doc, 1):
                if progress_callback:
                    progress_callback(10 + (i / total_pages * 20), f"Extracting slide {i}/{total_pages}...")
                
                # Extract text blocks from the page
                text = page.get_text()
                
                # Attempt to identify title and content
                lines = text.split('\n')
                title = ""
                content = ""
                
                # Skip empty lines at the beginning
                clean_lines = [line for line in lines if line.strip()]
                
                if clean_lines:
                    # Consider the first non-empty line as title
                    title = clean_lines[0].strip()
                    # Join the rest as content
                    content = '\n'.join(clean_lines[1:]).strip()
                
                # Store slide info even if title or content is minimally populated
                slides_content.append({
                    'title': title,
                    'content': content,
                    'slide_num': i
                })
                
                # Debug info
                print(f"Extracted PDF Slide {i}:")
                print(f"  Title: {title}")
                print(f"  Content length: {len(content)}")
            
            doc.close()
        except ImportError:
            # Fallback to PyPDF2 if PyMuPDF is not available
            print("PyMuPDF not available, falling back to PyPDF2 (less accurate extraction)")
            with open(pdf_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                total_pages = len(reader.pages)
                
                for i in range(total_pages):
                    if progress_callback:
                        progress_callback(10 + (i / total_pages * 20), f"Extracting slide {i+1}/{total_pages}...")
                    
                    page = reader.pages[i]
                    text = page.extract_text()
                    
                    # Simple slide content extraction
                    lines = text.split('\n')
                    title = ""
                    content = ""
                    
                    if lines:
                        title = lines[0].strip()
                        content = '\n'.join(lines[1:]).strip()
                    
                    # Store slide info
                    slides_content.append({
                        'title': title,
                        'content': content,
                        'slide_num': i+1
                    })
                    
                    # Debug info
                    print(f"Extracted PDF Slide {i+1}:")
                    print(f"  Title: {title}")
                    print(f"  Content length: {len(content)}")
        
        return slides_content
    
    def clean_slide_content(self, slide):
        """Clean and prepare slide content for Claude API"""
        # Remove common header/footer patterns
        title = slide['title']
        content = slide['content']
        
        # Common patterns to ignore in titles (like "June 1, 1999 Vi Editor X")
        header_patterns = [
            r"\w+ \d+, \d{4} .+ \d+",  # Date format followed by title and number
            r"^\d+$",  # Just a number
            r"Slide \d+",  # "Slide X"
        ]
        
        for pattern in header_patterns:
            if re.match(pattern, title):
                # Try to extract a better title from content if possible
                content_lines = content.split('\n')
                if content_lines and content_lines[0].strip():
                    # Move first line of content to title
                    title = content_lines[0].strip()
                    content = '\n'.join(content_lines[1:]).strip()
                break
        
        # Return the cleaned slide content
        return {
            'title': title,
            'content': content,
            'slide_num': slide['slide_num']
        }
    
    def generate_flashcards_with_claude(self, slides_content, progress_callback=None):
        """Use Claude to generate flashcards from slide content"""
        all_cards = []
        total_slides = len(slides_content)
        processed_slides = 0
        
        for slide in slides_content:
            processed_slides += 1
            if progress_callback:
                progress_callback(30 + (processed_slides / total_slides * 50), 
                                  f"Generating flashcards for slide {slide['slide_num']}/{total_slides}...")
            
            # Clean and prepare slide content
            cleaned_slide = self.clean_slide_content(slide)
            
            # If slide has minimal content, try to create a basic card
            has_meaningful_content = len(cleaned_slide['title']) > 3 or len(cleaned_slide['content']) > 10
            
            if not has_meaningful_content:
                print(f"Skipping slide {cleaned_slide['slide_num']} - insufficient content")
                continue
                
            # Combine title and content for context
            full_text = f"Title: {cleaned_slide['title']}\n\nContent: {cleaned_slide['content']}"
            
            # Generate flashcards using Claude with retries
            max_retries = 3
            retry_count = 0
            success = False
            
            while retry_count < max_retries and not success:
                try:
                    cards = self._ask_claude_for_cards(full_text)
                    
                    # Add slide reference to each card
                    for card in cards:
                        card['slide'] = f"Slide {cleaned_slide['slide_num']}"
                        card['context'] = cleaned_slide['title']
                    
                    all_cards.extend(cards)
                    print(f"Generated {len(cards)} cards for slide {cleaned_slide['slide_num']}")
                    success = True
                    
                except Exception as e:
                    retry_count += 1
                    print(f"Error generating cards for slide {cleaned_slide['slide_num']} (attempt {retry_count}): {e}")
                    time.sleep(1)  # Brief pause before retry
            
            # Create a basic card if all Claude attempts failed
            if not success:
                print(f"Falling back to basic card for slide {cleaned_slide['slide_num']}")
                if cleaned_slide['title']:
                    question = f"Explain the concept of: {cleaned_slide['title']}"
                    answer = cleaned_slide['content'] if cleaned_slide['content'] else "Review the slide content."
                    
                    all_cards.append({
                        'question': question,
                        'answer': answer,
                        'slide': f"Slide {cleaned_slide['slide_num']}",
                        'context': "Auto-generated (Claude API failed)"
                    })
        
        return all_cards
    
    def _ask_claude_for_cards(self, slide_text):
        """Ask Claude to generate question-answer pairs from the slide text"""
        prompt = """
        Please analyze this slide content from an educational presentation and create 1-5 Anki flashcards based on the key concepts.
        
        For each important concept, create a question that tests understanding and a comprehensive answer.
        
        Slide content:
        {slide_text}
        
        Format your response as a JSON array of objects with 'question' and 'answer' keys.
        Example:
        [
            {{"question": "What is the capital of France?", "answer": "Paris"}},
            {{"question": "What is the formula for calculating area of a circle?", "answer": "A = πr²"}}
        ]
        
        Only output valid JSON that can be parsed with json.loads() in Python.
        
        If there's not enough meaningful content to create flashcards, return an empty array: []
        """
        
        # Using Claude API
        response = self.client.messages.create(
            model="claude-3-7-sonnet-20250219",
            max_tokens=1000,
            temperature=0.7,
            system="You create high-quality flashcards from educational content. Always respond with valid JSON.",
            messages=[
                {"role": "user", "content": prompt.format(slide_text=slide_text)}
            ]
        )
        
        # Extract the response text
        cards_text = response.content[0].text
        
        # Process the response to extract cards
        try:
            # Try to parse the response as JSON
            cards = json.loads(cards_text)
            return cards
        except json.JSONDecodeError:
            # Fallback: Extract JSON pattern from the response
            try:
                # Look for a JSON array pattern
                json_pattern = re.search(r'\[\s*\{.*\}\s*\]', cards_text, re.DOTALL)
                if json_pattern:
                    json_str = json_pattern.group(0)
                    # Fix common JSON formatting issues
                    json_str = re.sub(r'(\w+):', r'"\1":', json_str)  # Add quotes to keys
                    json_str = re.sub(r':\s*"([^"]*)"', r': "\1"', json_str)  # Fix value quotes
                    cards = json.loads(json_str)
                    return cards
            except (json.JSONDecodeError, AttributeError):
                pass
            
            # Last resort: manual extraction
            cards = []
            
            # Look for question/answer patterns
            qa_pairs = re.findall(r'["\']{0,1}question["\']{0,1}\s*:\s*["\'](.*?)["\']\s*,\s*["\']{0,1}answer["\']{0,1}\s*:\s*["\'](.*?)["\']', 
                                  cards_text, re.DOTALL)
            
            for q, a in qa_pairs:
                cards.append({"question": q.strip(), "answer": a.strip()})
            
            if not cards:
                # Try a more lenient pattern
                q_blocks = re.findall(r'question["\']?:[\s"\']*(.+?)[\s"\']*,', cards_text)
                a_blocks = re.findall(r'answer["\']?:[\s"\']*(.+?)[\s"\']*[},]', cards_text)
                
                for i in range(min(len(q_blocks), len(a_blocks))):
                    q = re.sub(r'["\']', '', q_blocks[i]).strip()
                    a = re.sub(r'["\']', '', a_blocks[i]).strip()
                    cards.append({"question": q, "answer": a})
            
            return cards if cards else [{"question": "Review this slide", "answer": slide_text}]
    
    def create_anki_deck(self, cards, deck_name):
        """Create Anki deck from extracted content with custom name"""
        deck_id = random.randrange(1 << 30, 1 << 31)
        deck = genanki.Deck(deck_id, deck_name)
        
        for card in cards:
            note = genanki.Note(
                model=self.model,
                fields=[
                    card['question'], 
                    card['answer'], 
                    card['slide'],
                    card.get('context', '')
                ]
            )
            deck.add_note(note)
        
        return deck
    
    def process_file(self, file_path, deck_name, progress_callback=None):
        """Process a presentation file and create AI-enhanced Anki cards"""
        # Determine file type
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # Extract content
        if progress_callback:
            progress_callback(10, "Extracting slides...")
            
        if file_extension == '.pptx':
            slides_content = self.extract_from_pptx(file_path, progress_callback)
        elif file_extension == '.pdf':
            slides_content = self.extract_from_pdf(file_path, progress_callback)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        # Generate cards using Claude
        if progress_callback:
            progress_callback(30, "Generating flashcards with Claude...")
        cards = self.generate_flashcards_with_claude(slides_content, progress_callback)
        
        # Create Anki deck
        if progress_callback:
            progress_callback(80, "Creating Anki deck...")
        deck = self.create_anki_deck(cards, deck_name)
        
        # Generate output path in Downloads folder
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        downloads_folder = os.path.join(os.path.expanduser("~"), "Downloads")
        output_path = os.path.join(downloads_folder, f"{base_name}_flashcards.apkg")
        
        # Save Anki package
        if progress_callback:
            progress_callback(90, f"Saving Anki package to {output_path}...")
        package = genanki.Package(deck)
        package.write_to_file(output_path)
        
        if progress_callback:
            progress_callback(100, f"Created {len(cards)} AI-enhanced flashcards successfully!")
        
        return len(cards), output_path


class AnkiConverterApp(tk.Tk):
    def __init__(self):
        super().__init__()
        
        # Insert your Claude API key here
        self.api_key = "INSERT_YOUR_API_KEY_HERE"
        
        self.title("Presentation to Anki Flashcards Converter")
        self.geometry("700x700")
        self.configure(padx=20, pady=20)
        
        self.setup_ui()
    
    def setup_ui(self):
        # Create a main frame
        main_frame = ttk.Frame(self)
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # Header
        header_label = ttk.Label(main_frame, text="Convert Presentations to Anki Flashcards", font=("Arial", 16, "bold"))
        header_label.pack(pady=10)
        
        # File selection frame
        file_frame = ttk.LabelFrame(main_frame, text="Presentation/PDF File")
        file_frame.pack(fill=tk.X, padx=10, pady=10)
        
        # File path entry
        self.file_path_var = tk.StringVar()
        file_path_entry = ttk.Entry(file_frame, textvariable=self.file_path_var, width=50)
        file_path_entry.pack(side=tk.LEFT, padx=5, pady=10, fill=tk.X, expand=True)
        
        # Browse button
        browse_button = ttk.Button(file_frame, text="Browse", command=self.browse_file)
        browse_button.pack(side=tk.RIGHT, padx=5, pady=10)
        
        # Deck name frame
        deck_frame = ttk.LabelFrame(main_frame, text="Anki Deck Settings")
        deck_frame.pack(fill=tk.X, padx=10, pady=10)
        
        # Deck name entry
        ttk.Label(deck_frame, text="Deck Name:").grid(row=0, column=0, padx=5, pady=10, sticky=tk.W)
        self.deck_name_var = tk.StringVar(value="AI-Enhanced Flashcards")
        deck_name_entry = ttk.Entry(deck_frame, textvariable=self.deck_name_var, width=40)
        deck_name_entry.grid(row=0, column=1, padx=5, pady=10, sticky=tk.W)
        
        # Debug mode checkbox
        self.debug_var = tk.BooleanVar(value=True)
        debug_check = ttk.Checkbutton(deck_frame, text="Show detailed processing output", variable=self.debug_var)
        debug_check.grid(row=1, column=0, columnspan=2, padx=5, pady=5, sticky=tk.W)
        
        # Output info label
        ttk.Label(deck_frame, text="Output Location:").grid(row=2, column=0, padx=5, pady=10, sticky=tk.W)
        output_info = ttk.Label(deck_frame, text="Files will be saved to your Downloads folder automatically")
        output_info.grid(row=2, column=1, padx=5, pady=10, sticky=tk.W)
        
        # File drop placeholder frame
        drop_frame = ttk.LabelFrame(main_frame, text="File Instructions")
        drop_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        # Instructions label
        instructions = """
        1. Click "Browse" to select your PowerPoint (.pptx) or PDF (.pdf) presentation
        2. Enter a name for your Anki deck
        3. Click "Convert to Anki" to generate flashcards
        4. The Anki package (.apkg) will be saved to your Downloads folder
        

        """
        
        instructions_label = ttk.Label(drop_frame, text=instructions, font=("Arial", 11), justify=tk.LEFT)
        instructions_label.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)
        
        # Progress bar
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(main_frame, variable=self.progress_var, maximum=100)
        self.progress_bar.pack(fill=tk.X, padx=10, pady=10)
        
        # Status label
        self.status_var = tk.StringVar(value="Ready to convert")
        status_label = ttk.Label(main_frame, textvariable=self.status_var, font=("Arial", 10))
        status_label.pack(pady=5)
        
        # Convert button
        convert_button = ttk.Button(main_frame, text="Convert to Anki", command=self.convert_to_anki)
        convert_button.pack(pady=10)
    
    def browse_file(self):
        file_path = filedialog.askopenfilename(
            title="Select Presentation",
            filetypes=[("Presentation files", "*.pptx;*.pdf"), ("All files", "*.*")]
        )
        if file_path:
            self.file_path_var.set(file_path)
            # Set default deck name based on input file
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            self.deck_name_var.set(f"{base_name} Flashcards")
    
    def update_progress(self, value, message):
        self.progress_var.set(value)
        self.status_var.set(message)
        self.update_idletasks()
    
    def convert_to_anki(self):
        file_path = self.file_path_var.get()
        deck_name = self.deck_name_var.get()
        
        if not file_path:
            messagebox.showerror("Error", "Please select a presentation file.")
            return
        
        if not deck_name:
            messagebox.showerror("Error", "Please enter a deck name.")
            return
        
        # Disable the convert button during conversion
        for widget in self.winfo_children():
            if isinstance(widget, ttk.Button):
                widget.configure(state=tk.DISABLED)
        
        # Reset progress bar
        self.progress_var.set(0)
        self.status_var.set("Starting conversion...")
        
        # Run the conversion in a separate thread to keep the UI responsive
        def run_conversion():
            try:
                # Initialize the converter
                converter = ClaudeEnhancedConverter(self.api_key)
                
                # Process the file
                num_cards, output_path = converter.process_file(
                    file_path, 
                    deck_name,
                    progress_callback=self.update_progress
                )
                
                # Show success message
                self.after(0, lambda: messagebox.showinfo(
                    "Conversion Complete", 
                    f"Successfully created {num_cards} flashcards in:\n{output_path}"
                ))
                
            except Exception as e:
                # Show error message
                self.after(0, lambda: messagebox.showerror("Error", str(e)))
                self.update_progress(0, "Conversion failed.")
            
            # Re-enable the convert button
            self.after(0, self.enable_buttons)
        
        # Start the conversion thread
        threading.Thread(target=run_conversion, daemon=True).start()
    
    def enable_buttons(self):
        for widget in self.winfo_children():
            if isinstance(widget, ttk.Button):
                widget.configure(state=tk.NORMAL)

if __name__ == "__main__":
    app = AnkiConverterApp()
    app.mainloop()